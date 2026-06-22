# -*- coding: utf-8 -*-
"""KCS AI Agents 아키텍처 소개 덱 생성기"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
NAVY   = RGBColor(0x1B, 0x24, 0x4E)
BLUE   = RGBColor(0x06, 0x5A, 0x82)
TEAL   = RGBColor(0x1C, 0x72, 0x93)
SEA    = RGBColor(0x2A, 0x9D, 0x8F)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF8)
CARD   = RGBColor(0xE9, 0xEF, 0xF4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x20, 0x33)
MUTED  = RGBColor(0x5A, 0x6B, 0x7B)
GOLD   = RGBColor(0xE2, 0x9A, 0x2C)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)

FONT = "맑은 고딕"
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s

def box(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def title(s, t, sub=None, dark=False):
    tc = WHITE if dark else INK
    mc = ICE if dark else MUTED
    text(s, 0.62, 0.78, 11.5, 0.9, [[(t, 30, tc, True)]])
    if sub:
        text(s, 0.64, 1.46, 11.5, 0.45, [[(sub, 14, mc, False)]])

def kicker(s, t, color=TEAL, x=0.64, y=0.42):
    text(s, x, y, 6, 0.35, [[(t.upper(), 12, color, True)]])

def circle_num(s, x, y, d, n, fill=TEAL, tcol=WHITE):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.shadow.inherit = False; c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background()
    tf = c.text_frame; tf.word_wrap=False
    p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(15); r.font.bold=True
    r.font.color.rgb=tcol; r.font.name=FONT
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    return c

# ============================================================== 1. TITLE
s = slide(NAVY)
# subtle band
box(s, 0, 5.55, 13.333, 1.95, fill=RGBColor(0x16,0x1E,0x42))
text(s, 0.9, 1.75, 11, 0.5, [[("관세청 AI 통합 분석 플랫폼", 16, GOLD, True)]])
text(s, 0.88, 2.25, 11.6, 1.6,
     [[("KCS AI Agents", 60, WHITE, True)]])
text(s, 0.9, 3.55, 11.5, 1.0,
     [[("LangGraph 기반 멀티 에이전트 · 3중 데이터 계층 · 노드형 플로우 빌더", 19, ICE, False)]])
text(s, 0.92, 5.95, 11, 0.9,
     [[("전체 패키지 · 아키텍처 · 데이터 흐름 · 기술 스택", 14, ICE, False)],
      [("Python 3.11 · LangChain/LangGraph · DuckDB · Neo4j · ChromaDB · Vanilla JS", 12.5, RGBColor(0x9F,0xB0,0xD0), False)]])

# ============================================================== 2. OVERVIEW
s = slide(WHITE)
kicker(s, "Overview")
title(s, "한눈에 보는 플랫폼", "관세조사·수사 업무를 위한 AI 에이전트 오케스트레이션 플랫폼")
stats = [("30+", "독립 AI 에이전트", TEAL), ("3", "데이터 계층 (관계·그래프·벡터)", BLUE),
         ("4", "조사 분석 도메인", SEA), ("3", "LLM 프로바이더 라우팅", GOLD)]
x = 0.64
for val, lab, col in stats:
    box(s, x, 1.75, 2.92, 1.55, fill=LIGHT)
    box(s, x, 1.75, 0.12, 1.55, fill=col, shape=MSO_SHAPE.RECTANGLE)  # left marker kept minimal-> remove? keep subtle
    text(s, x+0.28, 1.92, 2.6, 0.9, [[(val, 40, col, True)]])
    text(s, x+0.30, 2.72, 2.55, 0.5, [[(lab, 12.5, INK, True)]], line_spacing=0.95)
    x += 3.06
# narrative cards
items = [
    ("무엇인가", "수입신고·기업거래·우범자 반입 데이터를 AI 에이전트로 다각 분석하고 보고서를 자동 생성하는 웹 플랫폼"),
    ("어떻게", "사용자가 조사 대상과 시나리오를 구성하면 LangGraph StateGraph가 에이전트들을 순차 실행하여 결과를 누적"),
    ("핵심 가치", "DuckDB(수치)·Neo4j(관계망)·ChromaDB(의미검색) 3중 데이터로 한국어 RAG 기반 심층 분석 제공"),
]
y = 3.7
for h, d in items:
    box(s, 0.64, y, 12.05, 1.02, fill=WHITE, line=CARD, line_w=1.0)
    text(s, 0.9, y+0.13, 2.6, 0.8, [[(h, 16, TEAL, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.5, y+0.13, 9.0, 0.8, [[(d, 14, INK, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y += 1.16

# ============================================================== 3. ARCHITECTURE
s = slide(LIGHT)
kicker(s, "Architecture")
title(s, "4계층 아키텍처", "프론트엔드 → 서버 → 워크플로우 오케스트레이션 → 에이전트 → 데이터")
layers = [
    ("프론트엔드 (Vanilla JS SPA)", "pages · analysis(4도메인) · core · shared 컴포넌트", BLUE),
    ("HTTP 서버  ·  web_server.py", "ThreadingHTTPServer · JSON API · UTF-8", TEAL),
    ("워크플로우 오케스트레이션  ·  workflows.py", "LangGraph StateGraph · CustomsState 상태 누적", SEA),
    ("AI 에이전트 계층  ·  src/agents/*.py", "분석 → 검증 → 관계망 → 생성 (30+ 모듈, 레지스트리 기반)", NAVY),
]
y = 1.85
for i, (h, d, col) in enumerate(layers):
    box(s, 1.4, y, 10.5, 0.92, fill=col)
    text(s, 1.75, y+0.12, 9.9, 0.7,
         [[(h, 16.5, WHITE, True)],[(d, 12, ICE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.45), Inches(y+0.93), Inches(0.4), Inches(0.18))
        ar.shadow.inherit=False; ar.fill.solid(); ar.fill.fore_color.rgb=MUTED; ar.line.fill.background()
    y += 1.12
# data tier strip
box(s, 1.4, y+0.02, 10.5, 0.78, fill=WHITE, line=NAVY, line_w=1.25)
text(s, 1.75, y+0.04, 9.9, 0.72,
     [[("데이터 계층    ", 14, NAVY, True), ("DuckDB (OLAP)   ·   Neo4j (그래프)   ·   ChromaDB (벡터/RAG)", 13.5, INK, False)]],
     anchor=MSO_ANCHOR.MIDDLE)

# ============================================================== 4. TECH STACK
s = slide(WHITE)
kicker(s, "Tech Stack")
title(s, "기술 스택", "언어 · AI 프레임워크 · 데이터 · 프론트엔드 · 인프라")
groups = [
    ("언어 · 런타임", BLUE, ["Python 3.10–3.12", "Vanilla JavaScript (ES6+)", "stdlib HTTP 서버"]),
    ("AI · LLM", TEAL, ["LangChain + LangGraph", "OpenAI · Anthropic · Google", "ko-sroberta 임베딩"]),
    ("데이터", SEA, ["DuckDB (로컬 OLAP)", "Neo4j 5.26+ (그래프)", "ChromaDB (벡터)"]),
    ("프론트엔드", NAVY, ["모듈식 SPA", "Drawflow (노드 UI)", "Dagre (그래프 레이아웃)"]),
    ("문서 처리", GOLD, ["PyPDF4 (PDF)", "openpyxl (Excel)", "pandas"]),
    ("인프라", TEAL, ["Docker · Compose", "Neo4j 컨테이너", ".env / YAML 설정"]),
]
gx, gy = 0.64, 1.85
cw, ch, gap = 3.93, 1.62, 0.13
for i, (h, col, rows) in enumerate(groups):
    cx = gx + (i % 3) * (cw + gap)
    cy = gy + (i // 3) * (ch + 0.22)
    box(s, cx, cy, cw, ch, fill=LIGHT)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+0.28), Inches(cy+0.26), Inches(0.2), Inches(0.2))
    dot.shadow.inherit=False; dot.fill.solid(); dot.fill.fore_color.rgb=col; dot.line.fill.background()
    text(s, cx+0.62, cy+0.16, cw-0.8, 0.4, [[(h, 15, INK, True)]])
    body = [[("·  "+r, 12.5, MUTED, False)] for r in rows]
    text(s, cx+0.3, cy+0.66, cw-0.5, 0.9, body, line_spacing=1.05, space_after=2)

# ============================================================== 5. BACKEND MODULES
s = slide(WHITE)
kicker(s, "Backend")
title(s, "백엔드 핵심 모듈", "src/ — 진입점부터 데이터 헬퍼까지")
rows = [
    ("web_server.py", "HTTP API 진입점 · ThreadingHTTPServer · JSON 응답"),
    ("workflows.py", "LangGraph StateGraph 조립 · create_initial_state()"),
    ("agents/state.py", "CustomsState (TypedDict) — 단계별 결과 누적 상태"),
    ("agents/module_registry.py", "AGENT_MODULES · 에이전트 등록/조회 레지스트리"),
    ("agents/service_registry.py", "AI 서비스/RAG 타입 설정 · 기본 프롬프트"),
    ("config.py", "YAML 설정 로더 (CFG 싱글턴 · thresholds.yaml)"),
    ("llm.py", "LLM 프로바이더 라우팅 (OpenAI/Anthropic/Gemini)"),
    ("embeddings.py", "HuggingFace 한국어 임베딩 싱글턴 캐시"),
    ("neo4j_graph.py", "관계망 그래프 빌드 (build_*_network_graph)"),
    ("paths.py / encoding.py", "경로 상수 관리 · UTF-8 입출력 강제"),
]
y = 1.8
half = 5
for col_i in range(2):
    cx = 0.64 + col_i * 6.18
    for r in range(half):
        name, desc = rows[col_i*half + r]
        yy = y + r * 1.04
        box(s, cx, yy, 5.9, 0.92, fill=LIGHT)
        text(s, cx+0.25, yy+0.1, 5.5, 0.4, [[(name, 13.5, BLUE, True)]])
        text(s, cx+0.25, yy+0.46, 5.5, 0.42, [[(desc, 11.5, MUTED, False)]], line_spacing=0.95)

# ============================================================== 6. AGENTS
s = slide(NAVY)
kicker(s, "AI Agents", color=GOLD)
title(s, "AI 에이전트 계층", "30+ 독립 에이전트가 4단계 파이프라인으로 협업", dark=True)
phases = [
    ("분석", BLUE, ["agent_company", "agent_db", "agent_rag", "agent_audit_search"]),
    ("검증", TEAL, ["declaration_verify", "hs_code", "customs_value", "origin"]),
    ("관계망", SEA, ["agent_network", "abnormal_trade", "proceeds", "person_risk"]),
    ("생성", GOLD, ["report", "summary", "translate", "result_synthesis"]),
]
x = 0.64
for i, (h, col, ags) in enumerate(phases):
    box(s, x, 1.9, 2.92, 4.6, fill=RGBColor(0x22,0x2C,0x5C))
    box(s, x, 1.9, 2.92, 0.72, fill=col)
    text(s, x, 1.95, 2.92, 0.62, [[(f"{i+1}. {h}", 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = 2.85
    for a in ags:
        box(s, x+0.22, yy, 2.48, 0.66, fill=RGBColor(0x2C,0x37,0x6E))
        text(s, x+0.36, yy+0.04, 2.3, 0.58, [[(a, 12.5, ICE, False)]], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.8
    x += 3.06
text(s, 0.64, 6.7, 12, 0.5, [[("레지스트리 기반 등록으로 에이전트 추가·교체·테스트가 독립적 — 결과는 모두 CustomsState 에 누적됩니다.", 12.5, RGBColor(0x9F,0xB0,0xD0), False)]])

# ============================================================== 7. DATA TIER
s = slide(WHITE)
kicker(s, "Data")
title(s, "3중 데이터 계층", "수치 · 관계망 · 의미검색을 각각 최적의 저장소로")
tiers = [
    ("DuckDB", "관계형 OLAP", BLUE, ["customs.duckdb (단일 파일)", "company · import · risk", "수치/집계 분석", "pandas .df() 연동"]),
    ("Neo4j", "그래프 DB", SEA, ["인물·기업·사건 노드", "거래/관계 엣지", "관계망 시각화 소스", "Docker 컨테이너"]),
    ("ChromaDB", "벡터 / RAG", GOLD, ["chroma_db/ (5개 컬렉션)", "ko-sroberta 임베딩", "의미 기반 검색", "한국어 문맥 이해"]),
]
x = 0.64
for name, sub, col, rows in tiers:
    box(s, x, 1.85, 3.93, 4.55, fill=LIGHT)
    box(s, x, 1.85, 3.93, 1.05, fill=col)
    text(s, x+0.3, 1.95, 3.5, 0.5, [[(name, 22, WHITE, True)]])
    text(s, x+0.32, 2.5, 3.5, 0.35, [[(sub, 12.5, WHITE, False)]])
    yy = 3.2
    for r in rows:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.32), Inches(yy+0.08), Inches(0.1), Inches(0.1))
        dot.shadow.inherit=False; dot.fill.solid(); dot.fill.fore_color.rgb=col; dot.line.fill.background()
        text(s, x+0.55, yy, 3.2, 0.6, [[(r, 13, INK, False)]])
        yy += 0.72
    x += 4.06

# ============================================================== 8. DATA FLOW
s = slide(LIGHT)
kicker(s, "Data Flow")
title(s, "데이터 흐름", "요청 한 번이 상태 그래프를 따라 흐르는 경로")
steps = [
    ("웹 UI", "조사대상·시나리오 입력", BLUE),
    ("web_server.py", "JSON 요청 파싱", TEAL),
    ("workflows.py", "CustomsState 초기화", SEA),
    ("에이전트 체인", "단계별 실행·결과 누적", NAVY),
    ("result_synthesis", "최종 종합 → JSON", GOLD),
    ("웹 UI 렌더", "대시보드·보고서 표시", BLUE),
]
x, y = 0.64, 2.2
bw, bh = 1.92, 1.5
for i, (h, d, col) in enumerate(steps):
    box(s, x, y, bw, bh, fill=WHITE, line=col, line_w=1.5)
    circle_num(s, x+0.18, y+0.16, 0.42, i+1, fill=col)
    text(s, x+0.15, y+0.68, bw-0.3, 0.4, [[(h, 13, INK, True)]], align=PP_ALIGN.CENTER, line_spacing=0.9)
    text(s, x+0.12, y+1.04, bw-0.24, 0.42, [[(d, 10.5, MUTED, False)]], align=PP_ALIGN.CENTER, line_spacing=0.9)
    if i < 5:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+bw+0.0), Inches(y+bh/2-0.1), Inches(0.16), Inches(0.2))
        ar.shadow.inherit=False; ar.fill.solid(); ar.fill.fore_color.rgb=MUTED; ar.line.fill.background()
    x += bw + 0.16
box(s, 0.64, 4.3, 12.05, 1.75, fill=WHITE, line=CARD, line_w=1.0)
text(s, 0.95, 4.5, 11.4, 1.4,
     [[("상태 기반 워크플로우 (StateGraph)", 15, TEAL, True)],
      [("각 에이전트의 출력은 공유 CustomsState 에 키별로 저장되어 다음 단계가 이전 결과를 참조합니다.", 13, INK, False)],
      [("사용자 진행 상태는 data/workspace_state/{userId}.json 에 개인별로 분리 저장 — 멀티 사용자 이력 관리.", 13, INK, False)]],
     line_spacing=1.15, space_after=6)

# ============================================================== 9. FRONTEND
s = slide(WHITE)
kicker(s, "Frontend")
title(s, "프론트엔드 구조", "web/static/js — 모듈식 단일 페이지 앱")
cols = [
    ("pages/", BLUE, ["home.js — 홈 컴포저", "scenario-builder.js", "agentic-service.js", "agentic-flow.js (Drawflow)", "intl.js — 국제정보"]),
    ("analysis/", SEA, ["customs/ — 관세조사", "general-investigation/", "special-investigation/", "shared/ — 공유 컴포넌트"]),
    ("core / shared", NAVY, ["dom.js — 렌더/마크다운", "page-registry.js — 라우팅", "tabs.js — 탭/서브탭", "network-graph.js", "prompt-composer.js"]),
]
x = 0.64
for h, col, rows in cols:
    box(s, x, 1.85, 3.93, 4.75, fill=LIGHT)
    text(s, x+0.3, 2.02, 3.5, 0.5, [[(h, 19, col, True)]])
    yy = 2.72
    for r in rows:
        box(s, x+0.28, yy, 3.4, 0.58, fill=WHITE)
        text(s, x+0.45, yy, 3.2, 0.58, [[(r, 12.5, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.7
    x += 4.06

# ============================================================== 10. DOMAINS
s = slide(WHITE)
kicker(s, "Domains")
title(s, "4가지 조사 분석 도메인", "도메인마다 전용 탭·상태·보고서 구성")
doms = [
    ("관세조사 (customs)", "기업 · company_id", BLUE, "수입신고 위험 · 원산지 · 과세가격", "dashboard · profile · risk · scenario · report"),
    ("기업수사 (general)", "기업 · company_id", TEAL, "거래 구조 · 부정 거래 · 관계사", "workbench · cases · profile · report"),
    ("우범자수사 (special)", "개인 · person_id", SEA, "반입 기록 · 동행자 · 여행경로", "forensic · network · profile · scenario"),
    ("국제정보 (intl)", "통합 정보", GOLD, "해외 정보 · 법령/특허 연계", "검색 · 요약 · 표준 보고서"),
]
gx, gy = 0.64, 1.85
cw, ch = 6.0, 2.18
for i, (name, tgt, col, focus, tabs) in enumerate(doms):
    cx = gx + (i % 2) * (cw + 0.2)
    cy = gy + (i // 2) * (ch + 0.2)
    box(s, cx, cy, cw, ch, fill=LIGHT)
    box(s, cx, cy, 0.16, ch, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, cx+0.4, cy+0.2, cw-0.7, 0.5, [[(name, 17, INK, True)]])
    text(s, cx+0.4, cy+0.72, cw-0.7, 0.35, [[("대상: ", 12, MUTED, True),(tgt, 12, col, True)]])
    text(s, cx+0.4, cy+1.12, cw-0.7, 0.4, [[("초점  ", 11.5, MUTED, True),(focus, 12.5, INK, False)]], line_spacing=0.95)
    text(s, cx+0.4, cy+1.6, cw-0.7, 0.4, [[("탭  ", 11.5, MUTED, True),(tabs, 11.5, MUTED, False)]], line_spacing=0.95)

# ============================================================== 11. KEY FEATURES
s = slide(LIGHT)
kicker(s, "Features")
title(s, "주요 기능 모듈", "분석 구성부터 노드형 플로우까지")
feats = [
    ("시나리오 빌더", "분석 단계를 수동 구성·저장·불러오기", "scenario-builder.js · scenario_templates.json", BLUE),
    ("AI 코칭 · 프롬프트 컴포저", "프롬프트 동적 생성·편집, 카드 코칭", "prompt-composer.js · prompt_overrides.json", TEAL),
    ("AI Agentic 플로우 빌더", "노드를 시각적으로 연결하는 파이프라인", "agentic-flow.js (Drawflow) · agentic-service.js", SEA),
    ("관계망 분석", "인물/기업 네트워크 시각화·시나리오", "network-graph.js · neo4j_graph.py", NAVY),
    ("한국어 RAG", "의미 기반 지식 검색", "agent_rag.py · embeddings.py · chroma_db", GOLD),
    ("워크스페이스 상태", "사용자별 진행 이력 분리 저장", "workspace_state/{userId}.json", TEAL),
]
gx, gy = 0.64, 1.85
cw, ch = 3.93, 2.18
for i, (h, d, path, col) in enumerate(feats):
    cx = gx + (i % 3) * (cw + 0.13)
    cy = gy + (i // 3) * (ch + 0.2)
    box(s, cx, cy, cw, ch, fill=WHITE, line=CARD, line_w=1.0)
    circle_num(s, cx+0.28, cy+0.26, 0.5, i+1, fill=col)
    text(s, cx+0.95, cy+0.28, cw-1.1, 0.7, [[(h, 14.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)
    text(s, cx+0.3, cy+1.05, cw-0.55, 0.5, [[(d, 12.5, MUTED, False)]], line_spacing=1.0)
    text(s, cx+0.3, cy+1.62, cw-0.55, 0.5, [[(path, 10.5, col, False)]], line_spacing=0.95)

# ============================================================== 12. DIRECTORY
s = slide(NAVY)
kicker(s, "Layout", color=GOLD)
title(s, "디렉터리 구조", "최상위 패키지 한눈에", dark=True)
tree = [
    ("src/", "Python 백엔드 — agents/ · workflows · llm · embeddings · neo4j_graph"),
    ("web/static/", "프론트엔드 — js/{pages,analysis,core} · prompts/ · vendor/"),
    ("data/", "customs.duckdb · chroma_db/ · workspace_state/ · *.json 설정"),
    ("config/", "thresholds.yaml — 위험지표 임계값"),
    ("files/", "사용자 업로드 (Import_Declaration 등)"),
    ("scripts/ · docs/", "데이터 시드 스크립트 · 문서"),
    ("web_server.py", "HTTP 서버 진입점"),
    ("requirements.txt · docker-compose*.yml · .env.example", "의존성 · 컨테이너 · 환경설정"),
]
y = 1.95
for name, desc in tree:
    box(s, 0.64, y, 12.05, 0.56, fill=RGBColor(0x22,0x2C,0x5C))
    text(s, 0.9, y+0.02, 4.6, 0.52, [[(name, 13.5, GOLD, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.4, y+0.02, 7.1, 0.52, [[(desc, 12, ICE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.64

# ============================================================== 13. TIMELINE
s = slide(WHITE)
kicker(s, "Evolution")
title(s, "최근 진화", "커밋으로 본 기능 발전 방향")
commits = [
    ("6f2f310", "AI Agentic 서비스 — 노드 기반 플로우 빌더 추가", SEA),
    ("358d9ff", "프레임 동적 리사이즈 (드래그 거터)", TEAL),
    ("94ac72b", "입력 모델 단일화 · 카드 AI 코칭", BLUE),
    ("51fac52", "관계망 분석 시나리오 저장·적용", SEA),
    ("da1592a", "홈 컴포저 가로 흐름 · 프롬프트 자동 생성", TEAL),
    ("cd53167", "진행 작업 상태 사용자별 분리", BLUE),
]
# vertical line
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.95), Inches(0.04), Inches(4.5))
ln.shadow.inherit=False; ln.fill.solid(); ln.fill.fore_color.rgb=CARD; ln.line.fill.background()
y = 1.95
for sha, msg, col in commits:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.36), Inches(y+0.12), Inches(0.32), Inches(0.32))
    dot.shadow.inherit=False; dot.fill.solid(); dot.fill.fore_color.rgb=col; dot.line.color.rgb=WHITE; dot.line.width=Pt(2)
    box(s, 2.0, y, 10.65, 0.62, fill=LIGHT)
    text(s, 2.25, y+0.02, 1.5, 0.58, [[(sha, 12.5, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, y+0.02, 8.8, 0.58, [[(msg, 13, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.74

# ============================================================== 14. CLOSING
s = slide(NAVY)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
text(s, 0.9, 1.3, 11, 0.5, [[("Summary", 14, GOLD, True)]])
text(s, 0.88, 1.75, 11.5, 1.0, [[("핵심 요약", 40, WHITE, True)]])
points = [
    ("모듈식 멀티 에이전트", "30+ 독립 에이전트 · 레지스트리 기반 · StateGraph 오케스트레이션"),
    ("3중 데이터 계층", "DuckDB(수치) · Neo4j(관계망) · ChromaDB(한국어 RAG)"),
    ("4개 조사 도메인", "관세조사 · 기업수사 · 우범자수사 · 국제정보"),
    ("노드형 확장성", "Drawflow 플로우 빌더 · 프롬프트 컴포저로 비개발자도 구성"),
]
y = 3.05
for h, d in points:
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y+0.08), Inches(0.22), Inches(0.22))
    dot.shadow.inherit=False; dot.fill.solid(); dot.fill.fore_color.rgb=GOLD; dot.line.fill.background()
    text(s, 1.35, y, 11, 0.45, [[(h+"   ", 17, WHITE, True),(d, 13.5, ICE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.82
text(s, 0.92, 6.7, 11, 0.5, [[("KCS AI Agents · 관세청 AI 통합 분석 플랫폼", 12.5, RGBColor(0x9F,0xB0,0xD0), False)]])

prs.save("KCS_AI_Agents_Architecture.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
